"""
Rayuela — motor de plantillas PPTX.

NO genera PowerPoints desde cero: abre una COPIA del molde oficial del ICBF
(la misma plantilla que Jimena ya usa y tiene aprobada) y reemplaza únicamente
el texto dentro de los contenedores de respuesta correctos, dejando intactos
diseño, fuentes, colores, logos y posiciones.

Esta es exactamente la técnica que validamos y usamos para corregir y
completar los 34 cuadernos reales de Jimena (ver find_question /
find_answer_table / write_answer del proyecto original).
"""

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

TAG = ""
TEXT_COLOR = RGBColor(0x00, 0x00, 0x00)


def _normalizar(texto):
    """Convierte cualquier secuencia de espacios/saltos de línea en un solo
    espacio. Algunas preguntas en los moldes oficiales del ICBF traen un
    salto de línea en medio (p. ej. "...la intencionalidad del\nencuentro?"),
    mientras que en el código las escribimos con un espacio normal — sin esto,
    la búsqueda por subcadena no las encuentra y la casilla queda vacía."""
    return " ".join(texto.split())


LARGO_MAXIMO_ETIQUETA = 285
"""Las etiquetas/preguntas reales de los moldes del ICBF pueden llegar a
este tamaño. La más larga encontrada es la primera pregunta del 'Registro de
observaciones al desarrollo infantil mensual' del molde de llamada
('¿Qué le gustó a (nombre de la niña o niño) del encuentro o acompañamiento?
¿Qué no le gustó? ¿A que jugó? ¿Sobre qué conversó? ...' — 279 caracteres);
se eligió 285 para tener un pequeño margen por encima de ella.

Los párrafos instructivos/de advertencia que traen los moldes (los que no son
etiquetas de preguntas sino texto de guía para Jimena) miden 300 caracteres
o más (306, 473...) — por eso este umbral los sigue excluyendo correctamente.

HISTORIA DEL VALOR:
- Originalmente era 200, basado en la etiqueta más larga conocida en ese
  momento (124 chars, la de "paquete didáctico"). Sin este límite,
  find_question encontraba el párrafo instructivo de la diapositiva de
  INTENCIONALIDAD ('...que la intencionalidad debe estar ligada...', ~473
  chars) ANTES de llegar a la celda real, y la respuesta terminaba escrita
  en el lugar equivocado (Bug 1 reportado por Jimena: "la intencionalidad
  está mal organizada, la deja arriba y va abajo").
- Se subió a 285 al descubrir que la primera pregunta del 'Registro de
  observaciones' (slide 7 del molde de llamada) tiene 279 chars y quedaba
  excluida con el umbral de 200, dejando esa casilla vacía (Bug 6)."""


def find_question(slide, contains):
    """Busca el cuadro de texto de una pregunta por una subcadena de su texto
    (ignorando saltos de línea y espacios extra, que varían entre moldes).

    Ignora formas cuyo texto completo sea más largo que LARGO_MAXIMO_ETIQUETA
    (ver su docstring) para no confundir un párrafo instructivo largo con la
    etiqueta real de la pregunta."""
    contains_norm = _normalizar(contains).lower()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        texto_normalizado = _normalizar(shape.text_frame.text)
        if len(texto_normalizado) > LARGO_MAXIMO_ETIQUETA:
            continue
        if contains_norm in texto_normalizado.lower():
            return shape
    return None


def _bbox(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _contenedor_mas_cercano(slide, bbox_referencia, excluir_ids, tol=0.35):
    """
    Núcleo común de búsqueda espacial: dado el rectángulo (bbox) de una
    pregunta —ya sea un cuadro de texto suelto o una celda de tabla— busca
    entre las formas de la diapositiva el contenedor de respuesta (tabla,
    freeform o agrupado) más cercano y alineado con ella.

    OJO: usamos el borde SUPERIOR de la pregunta (y0) como referencia, no el
    inferior — muchos cuadros de texto de preguntas en los moldes del ICBF
    tienen una altura "de diseño" mucho mayor que el texto real (quedan con
    bastante espacio vacío debajo), así que su borde inferior puede estar más
    cerca de la casilla de la SIGUIENTE pregunta que de la suya propia. Eso
    fue justo lo que causaba que algunas respuestas se escribieran en la
    casilla equivocada (y otras se sobreescribieran entre sí).

    Priorizamos la alineación horizontal (misma columna) sobre la cercanía
    vertical, y solo consideramos contenedores que empiezan en o debajo del
    inicio de la pregunta.
    """
    qx0, qy0, qx1, qy1 = bbox_referencia
    margen = Pt(tol * 72)
    candidates = []
    for shape in slide.shapes:
        if shape.shape_id in excluir_ids:
            continue
        if shape.shape_type in (19, 5, 6):  # TABLE, FREEFORM, GROUP
            # Las tablas SIEMPRE pueden recibir texto (celda 0,0). Las formas
            # sueltas o agrupadas, en cambio, solo sirven como contenedor si
            # tienen su propio cuadro de texto — de lo contrario son elementos
            # decorativos (líneas, bordes, marcos) que por su posición pueden
            # "ganarle" por muy poco a la casilla real y dejarla vacía. Esto
            # fue justo lo que pasó con un grupo decorativo justo encima de la
            # tabla de respuesta de "Aspectos a tener en cuenta..." en el
            # molde de llamada — sin este filtro, la respuesta se perdía.
            if shape.shape_type != 19 and not shape.has_text_frame:
                continue
            sx0, sy0, sx1, sy1 = _bbox(shape)
            if sy0 < qy0 - margen:
                continue  # está por encima de la pregunta: no puede ser su respuesta
            dist = abs(sx0 - qx0) * 3 + abs(sy0 - qy0)
            candidates.append((dist, shape))
    if not candidates:
        return None
    # SEGUNDO FILTRO (descubierto al revisar "Registro de observaciones al
    # desarrollo infantil mensual"): incluso teniendo su propio cuadro de
    # texto, hay formas FREEFORM/GROUP puramente decorativas — el "marco" o
    # fondo dibujado casi exactamente encima de la tabla de respuesta real —
    # que por estar un poco más cerca de la pregunta le ganan el primer lugar
    # y dejan la tabla verdadera vacía. En los 34 cuadernos reales y en TODAS
    # las demás secciones de ambos moldes (familia, talento humano), el
    # contenedor de respuesta SIEMPRE terminó siendo una TABLA — nunca una
    # forma suelta o agrupada. Por eso, si entre los candidatos hay alguna
    # tabla, la preferimos sobre cualquier forma decorativa, y solo usamos
    # una forma suelta como respaldo si de verdad no hay ninguna tabla cerca.
    candidates.sort(key=lambda c: (0 if c[1].shape_type == 19 else 1, c[0]))
    return candidates[0][1]


def find_answer_table(slide, qshape, tol=0.35):
    """
    Encuentra el contenedor de respuesta espacialmente más cercano/alineado
    con la pregunta dada, cuando la pregunta es un cuadro de texto suelto.
    Delega el cálculo a `_contenedor_mas_cercano` (ver su docstring).
    """
    if qshape is None:
        return None
    return _contenedor_mas_cercano(slide, _bbox(qshape), {qshape.shape_id}, tol=tol)


def _bbox_celda(tabla_shape, fila, columna):
    """Calcula el rectángulo (x0,y0,x1,y1) de una celda dentro de una tabla,
    sumando los anchos/altos de las columnas/filas anteriores — python-pptx
    no expone la posición de una celda directamente, solo la de la tabla
    completa y el tamaño de cada columna/fila."""
    tabla = tabla_shape.table
    x0 = tabla_shape.left + sum(tabla.columns[c].width for c in range(columna))
    y0 = tabla_shape.top + sum(tabla.rows[r].height for r in range(fila))
    ancho = tabla.columns[columna].width
    alto = tabla.rows[fila].height
    return (x0, y0, x0 + ancho, y0 + alto)


def find_answer_table_para_celda(slide, tabla_shape, fila, columna, tol=0.35):
    """
    Encuentra el contenedor de respuesta para una pregunta que está escrita
    DENTRO de una celda de tabla (ver find_question_en_tabla).

    Al revisar la página de "voces del talento humano" del molde de llamada
    descubrimos que, para estas 2 preguntas en particular, la "celda vecina"
    (misma fila, columna siguiente) NO es el área de respuesta real: es apenas
    el resto visual de una celda combinada (columnas de 0.17 y 0.26 pulgadas
    de ancho — imposible que quepa una respuesta ahí sin desbordarse, que fue
    justo el problema que vio Alexander). La respuesta real va en una TABLA
    aparte, ubicada espacialmente cerca de la celda — exactamente el mismo
    patrón "pregunta en cuadro de texto + tabla de respuesta cercana" que usan
    las otras preguntas de esa misma página. Por eso reutilizamos la misma
    búsqueda espacial (`_contenedor_mas_cercano`), partiendo del rectángulo de
    la celda en lugar del rectángulo de un cuadro de texto.
    """
    bbox_pregunta = _bbox_celda(tabla_shape, fila, columna)
    return _contenedor_mas_cercano(slide, bbox_pregunta, {tabla_shape.shape_id}, tol=tol)


def find_shape_by_id(slide_or_group, shape_id):
    for shape in slide_or_group.shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == 6:  # GROUP
            found = find_shape_by_id(shape, shape_id)
            if found is not None:
                return found
    return None


def _escribir_en_textframe(tf, text, size):
    """Limpia un text_frame y escribe `text` (precedido del TAG) en su único
    párrafo, con el tamaño y color indicados. Lógica compartida entre
    write_answer (contenedor normal) y write_answer_en_celda (pregunta y
    respuesta dentro de la misma tabla)."""
    tf.word_wrap = True
    for p in list(tf.paragraphs[1:] if len(tf.paragraphs) > 1 else []):
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    r = p.add_run()
    r.text = TAG + text
    r.font.size = Pt(size)
    r.font.italic = False
    r.font.color.rgb = TEXT_COLOR


def write_answer(container, text, size=12):
    """Escribe `text` (precedido del TAG de borrador) dentro del contenedor."""
    if container is None:
        return False
    if container.shape_type == 19:  # TABLE
        tf = container.table.cell(0, 0).text_frame
    else:
        if not container.has_text_frame:
            return False
        tf = container.text_frame
        # Antes usábamos NONE (tamaño fijo). Pero los textos de "voces del
        # talento humano" suelen ser más largos que los de la familia y el
        # cuadro de esa diapositiva es más pequeño: con tamaño fijo el texto
        # se desbordaba y se montaba sobre las casillas vecinas (lo que vio
        # Alexander en el pantallazo). TEXT_TO_FIT_SHAPE le dice a PowerPoint
        # que reduzca automáticamente el tamaño de letra hasta que el texto
        # quepa dentro del cuadro — así nunca se desborda, sin importar cuánto
        # redacte la IA.
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _escribir_en_textframe(tf, text, size)
    return True


def find_question_en_tabla(slide, contains):
    """Busca una pregunta escrita DENTRO de una celda de tabla, en vez de en
    un cuadro de texto suelto (que es lo normal en los demás moldes).

    Esto pasa, por ejemplo, en la diapositiva de "voces del talento humano"
    del molde de llamada: dos de las cuatro preguntas quedaron escritas una
    junto a otra dentro de la MISMA tabla, con su casilla de respuesta vacía
    justo al lado (misma fila, columna siguiente) — un diseño distinto al de
    "pregunta en cuadro de texto + tabla de respuesta separada" que usa
    find_question/find_answer_table. Si no se busca también aquí, esas
    preguntas son invisibles para Rayuela y sus casillas quedan vacías.

    Devuelve (tabla_shape, fila, columna_de_la_pregunta) o None.
    """
    contains_norm = _normalizar(contains).lower()
    for shape in slide.shapes:
        if shape.shape_type != 19:  # TABLE
            continue
        tabla = shape.table
        for r, row in enumerate(tabla.rows):
            for c, cell in enumerate(row.cells):
                if contains_norm in _normalizar(cell.text_frame.text).lower():
                    return (shape, r, c)
    return None


def llenar_respuestas(slide, mapa_pregunta_a_texto, tol=0.35, size=12):
    """
    Recorre un mapa {subcadena_de_pregunta: texto_redactado} y escribe cada
    respuesta en el contenedor correspondiente de la diapositiva.

    Devuelve un reporte {pregunta: True/False} para registrar éxito/fallo
    (esto es justo lo que nos permitió detectar y corregir los problemas
    reales en el proyecto de los 34 cuadernos).
    """
    reporte = {}
    for pregunta, texto in mapa_pregunta_a_texto.items():
        q = find_question(slide, pregunta)
        if q is not None:
            contenedor = find_answer_table(slide, q, tol=tol)
            ok = write_answer(contenedor, texto, size=size)
        else:
            # No apareció como cuadro de texto suelto — puede que esté escrita
            # dentro de una celda de tabla (ver find_question_en_tabla).
            ubicacion = find_question_en_tabla(slide, pregunta)
            if ubicacion is not None:
                tabla_shape, fila, col = ubicacion
                contenedor = find_answer_table_para_celda(slide, tabla_shape, fila, col, tol=tol)
                ok = write_answer(contenedor, texto, size=size)
            else:
                ok = False
        reporte[pregunta] = ok
    return reporte
