"""
Rayuela — backend principal (Flask).

Flujo que implementa (módulo "voces", el primero — ya validado a fondo con
los 34 cuadernos reales de Jimena):

  1. Jimena manda: nombre del niño, fecha de nacimiento, sexo, tipo de
     cuaderno (hogar/llamada), actividad realizada, y sus observaciones
     en bruto.
  2. El backend calcula la banda de edad (no depende de que Jimena la
     recuerde — elimina justo el tipo de error que detectamos: "a Hassan
     le puso bebé y es un niño de un año").
  3. Llama a la API de Claude (capa `redactor`) para redactar cada
     respuesta en primera persona familiar, con el vocabulario correcto
     para esa banda de edad.
  4. Inserta el texto en una copia del molde .pptx oficial (capa
     `plantilla_pptx`) y entrega el archivo final.

Variables de entorno esperadas (configurar en Railway):
  ANTHROPIC_API_KEY   -> la llave de la API (la pone Alexander)
  RAYUELA_TEMPLATE_DIR -> carpeta donde están los moldes .pptx oficiales
"""

import os
import io
import uuid
import concurrent.futures
from datetime import datetime

from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation

from . import edades
from . import redactor
from . import plantilla_pptx
from . import db

app = Flask(__name__)
CORS(app)  # permite que el frontend (Netlify) llame a este backend (Railway)

# Crea la tabla planeaciones si no existe todavía — se ejecuta una vez
# al arrancar el servidor, sin tocar datos existentes.
with app.app_context():
    db.init()

TEMPLATE_DIR = os.environ.get("RAYUELA_TEMPLATE_DIR", os.path.join(
    os.path.dirname(__file__), "..", "templates"))

# Preguntas de "voces de la familia" tal como aparecen en los moldes reales del ICBF
PREGUNTAS_HOGAR = [
    "Qué les gustó y que no del encuentro",
    "sugerencias tiene la familia",
    "Lo que aprendieron de este encuentro",
    "Compromisos para el próximo encuentro",
]

PREGUNTAS_LLAMADA = [
    "Qué les gustó y que no de los acompañamientos a distancia",
    "Para qué le sirvieron los acompañamientos a distancia",
    "Cómo participó la niña, el niño o mujer gestante",
    "Qué le gustaría vivir en los próximos acompañamientos a distancia",
]

SLIDE_FAMILIA = {"hogar": 2, "llamada": 4}

# Preguntas de "voces del talento humano del servicio" (la reflexión de Jimena
# como agente educativa, NO la voz de la familia). En el molde de hogar son 5
# preguntas; en el de llamada el molde oficial solo trae 2.
PREGUNTAS_TALENTO_HOGAR = [
    "Considera que se alcanzó la intencionalidad del encuentro",
    "Cuál fue el mejor momento del encuentro",
    "Cómo participó la familia",
    "Qué recomendaciones harían para el próximo encuentro",
    "Cómo se aprovecharon los materiales propuestos",
]

PREGUNTAS_TALENTO_LLAMADA = [
    "Se cumplieron las intencionalidades propuestas",
    "Describa cómo fue la participación de la familia en el acompañamiento",
    "Cómo se vinculó la niña, el niño o mujer gestantes en los acompañamientos a distancia",
    "Qué recomendaciones tienen para los próximos acompañamientos a distancia",
]

SLIDE_TALENTO = {"hogar": 3, "llamada": 5}

# El molde de llamada trae, además, una sección que el de hogar NO tiene:
# "Registro de observaciones al desarrollo infantil mensual" — un registro
# narrativo, en tercera persona, sobre cómo le fue al niño/a la niña ese mes
# (gustos, participación, intereses, avances). Solo aplica a "llamada".
PREGUNTAS_DESARROLLO_LLAMADA = [
    "Qué le gustó a (nombre de la niña o niño) del encuentro o acompañamiento",
    "Qué intereses tiene",
    "Aspectos a tener en cuenta en los próximos encuentros o acompañamiento del siguiente mes",
]

SLIDE_DESARROLLO = {"llamada": 7}


@app.get("/api/salud")
def salud():
    return jsonify({"estado": "ok", "app": "Rayuela", "hora": datetime.now().isoformat()})


@app.post("/api/calcular-edad")
def calcular_edad():
    """Utilidad para que el frontend muestre la banda de edad en vivo, sin que
    Jimena tenga que calcularla ni recordarla."""
    datos = request.get_json(force=True)
    try:
        nacimiento = datetime.strptime(datos["fecha_nacimiento"], "%Y-%m-%d").date()
    except (KeyError, ValueError):
        return jsonify({"error": "fecha_nacimiento inválida, use AAAA-MM-DD"}), 400

    meses = edades.edad_en_meses(nacimiento)
    banda = edades.banda_por_edad(meses)
    genero = datos.get("genero", "M")
    sust = edades.sustantivo(genero, banda["clave"])
    return jsonify({
        "edad_meses": meses,
        "edad_legible": edades.formato_legible(meses),
        "banda": banda,
        "sustantivo": sust,
    })


@app.post("/api/generar-voces")
def generar_voces():
    """
    Genera el documento de 'voces de familia' ya diligenciado.

    Cuerpo esperado (JSON):
    {
      "nombre": "Hassan Melo Hueso",
      "fecha_nacimiento": "2024-09-09",   (no requerido si tipo_participante=="gestante")
      "genero": "M",                       (no requerido si tipo_participante=="gestante")
      "tipo_cuaderno": "hogar" | "llamada",
      "tipo_participante": "nino" | "nina" | "bebe" | "gestante"   (opcional, solo llamada)
      "actividad": "exploración sensorial con texturas",
      "observaciones": "le gustó tocar las telas, se rio mucho, ..."
    }
    """
    datos = request.get_json(force=True)
    tipo = (datos.get("tipo_cuaderno") or "").strip().lower()
    tipo_participante = (datos.get("tipo_participante") or "").strip().lower()
    es_gestante = tipo == "llamada" and tipo_participante == "gestante"

    requeridos = ["nombre", "tipo_cuaderno", "actividad", "observaciones"]
    if not es_gestante:
        requeridos += ["fecha_nacimiento", "genero"]
    faltantes = [c for c in requeridos if not datos.get(c)]
    if faltantes:
        return jsonify({"error": f"faltan campos: {', '.join(faltantes)}"}), 400

    if tipo not in SLIDE_FAMILIA:
        return jsonify({"error": "tipo_cuaderno debe ser 'hogar' o 'llamada'"}), 400

    # 1. Calculamos la banda de edad — sin depender de la memoria de Jimena
    #    (la mujer gestante no tiene fecha de nacimiento ni banda de edad)
    if es_gestante:
        banda = {"clave": "gestante", "etiqueta": "Mujer gestante"}
    else:
        try:
            nacimiento = datetime.strptime(datos["fecha_nacimiento"], "%Y-%m-%d").date()
        except ValueError:
            return jsonify({"error": "fecha_nacimiento inválida, use AAAA-MM-DD"}), 400
        meses = edades.edad_en_meses(nacimiento)
        banda = edades.banda_por_edad(meses)

    # 2. Redactamos cada respuesta con la API de Claude, en el tono correcto
    preguntas = PREGUNTAS_HOGAR if tipo == "hogar" else PREGUNTAS_LLAMADA
    materia_prima = (
        f"Actividad realizada: {datos['actividad']}\n"
        f"Lo que observé / lo que pasó: {datos['observaciones']}"
    )

    instrucciones = {
        preguntas[0]: "la respuesta de la familia en primera persona ('nosotros como "
                      "familia') sobre qué les gustó y qué no del encuentro/acompañamiento",
        preguntas[1]: "la respuesta de la familia en primera persona sobre sugerencias "
                      "para próximas experiencias",
        preguntas[2]: "la respuesta de la familia en primera persona sobre lo que "
                      "aprendieron de este encuentro/acompañamiento",
        preguntas[3]: "la respuesta de la familia en primera persona sobre los "
                      "compromisos para el próximo encuentro",
    }

    # 2b. Redactamos también las "voces del talento humano del servicio"
    # (la reflexión profesional de Jimena como agente educativa — tono
    # analítico, en primera persona singular, NO la voz de la familia).
    preguntas_talento = (PREGUNTAS_TALENTO_HOGAR if tipo == "hogar"
                         else PREGUNTAS_TALENTO_LLAMADA)

    instrucciones_talento = {
        preguntas_talento[0]: "mi reflexión como agente educativa, en primera persona "
                              "singular ('considero', 'observé'), sobre si se alcanzó "
                              "la intencionalidad pedagógica propuesta para este "
                              "encuentro/acompañamiento, con base en lo que ocurrió",
        preguntas_talento[1]: "mi reflexión como agente educativa, en primera persona "
                              "singular, sobre cuál fue el mejor momento del encuentro/"
                              "acompañamiento y por qué, desde mi mirada profesional",
    }
    if tipo == "hogar":
        instrucciones_talento.update({
            preguntas_talento[2]: "mi análisis como agente educativa, en primera "
                                  "persona singular, sobre cómo participó la familia "
                                  "y quiénes participaron en el encuentro",
            preguntas_talento[3]: "mis recomendaciones profesionales, en primera "
                                  "persona singular, para tener en cuenta en el "
                                  "próximo encuentro",
            preguntas_talento[4]: "mi valoración como agente educativa, en primera "
                                  "persona singular, sobre cómo se aprovecharon los "
                                  "materiales propuestos durante el encuentro",
        })
    else:
        instrucciones_talento[preguntas_talento[1]] = (
            "mi descripción y análisis, como agente educativa y en primera persona "
            "singular, de cómo fue la participación de la familia durante el "
            "acompañamiento a distancia"
        )
        instrucciones_talento[preguntas_talento[2]] = (
            "mi análisis, como agente educativa y en primera persona singular, "
            "sobre cómo se vinculó e involucró la niña, el niño o la mujer "
            "gestante protagonista en los acompañamientos a distancia realizados "
            "este mes"
        )
        instrucciones_talento[preguntas_talento[3]] = (
            "mis recomendaciones profesionales, en primera persona singular, "
            "para tener en cuenta en los próximos acompañamientos a distancia"
        )

    # 2c-bis. "Registro de observaciones al desarrollo infantil mensual" —
    # esta sección SOLO existe en el molde de llamada (el de hogar no la
    # trae). Es un registro narrativo en TERCERA PERSONA sobre cómo le fue
    # al niño/a la niña ese mes (gustos, participación, intereses, avances),
    # no la voz de la familia ni la reflexión de Jimena sobre su actividad.
    preguntas_desarrollo = (PREGUNTAS_DESARROLLO_LLAMADA
                            if (tipo == "llamada" and not es_gestante) else [])
    instrucciones_desarrollo = {}
    if tipo == "llamada" and not es_gestante:
        nombre_nino = datos["nombre"].strip().split()[0].title()
        instrucciones_desarrollo = {
            preguntas_desarrollo[0]: (
                f"una descripción narrativa, en tercera persona y nombrando a "
                f"{nombre_nino} por su nombre, sobre qué le gustó y qué no del "
                f"encuentro o acompañamiento, a qué jugó, sobre qué conversó, "
                f"cómo participó y cómo se relacionó con los adultos de la "
                f"familia, con el talento humano y con otras niñas y niños — "
                f"con base en lo que Jimena observó"
            ),
            preguntas_desarrollo[1]: (
                f"una descripción narrativa, en tercera persona y nombrando a "
                f"{nombre_nino} por su nombre, sobre qué intereses tiene, cómo "
                f"comunica sus intereses y necesidades, cómo se relaciona con su "
                f"familia, y qué aspectos de su desarrollo está apropiando y "
                f"comprendiendo — con base en lo que Jimena observó"
            ),
            preguntas_desarrollo[2]: (
                f"mis recomendaciones profesionales, en primera persona singular "
                f"como agente educativa, sobre los aspectos a tener en cuenta en "
                f"los próximos encuentros o acompañamientos del siguiente mes "
                f"para favorecer el proceso de desarrollo y aprendizaje de "
                f"{nombre_nino}"
            ),
        }

    # 2c. Lanzamos TODAS las llamadas a Claude EN PARALELO (antes se hacían
    # una por una, en fila — con 9 llamadas eso tardaba varios minutos). Cada
    # llamada tarda lo mismo, pero al lanzarlas todas a la vez el tiempo total
    # se acerca al de UNA sola llamada (unos 20-40 segundos) en vez de la suma
    # de las 9. Esto es lo que va a hacer que generar el documento sea mucho
    # más rápido para Jimena.
    trabajos = []
    for pregunta in preguntas:
        trabajos.append(("familia", pregunta, dict(
            banda_clave=banda["clave"],
            banda_etiqueta=banda["etiqueta"],
            instruccion=instrucciones[pregunta],
            materia_prima=materia_prima,
            evitar_actividad=(tipo == "llamada"),
        )))
    for pregunta in preguntas_talento:
        trabajos.append(("talento", pregunta, dict(
            banda_clave=banda["clave"],
            banda_etiqueta=banda["etiqueta"],
            instruccion=instrucciones_talento[pregunta],
            materia_prima=materia_prima,
            perspectiva="talento_humano",
            # El cuadro de "voces del talento humano" es más chico (caben 5
            # preguntas donde el de familia solo tiene 4): pedimos un texto
            # más corto para que quepa sin desbordarse.
            max_palabras=70,
            evitar_actividad=(tipo == "llamada"),
        )))
    for pregunta in preguntas_desarrollo:
        trabajos.append(("desarrollo", pregunta, dict(
            banda_clave=banda["clave"],
            banda_etiqueta=banda["etiqueta"],
            instruccion=instrucciones_desarrollo[pregunta],
            materia_prima=materia_prima,
            perspectiva="desarrollo_infantil",
            max_palabras=90,
            evitar_actividad=True,
        )))

    mapa_textos = {}
    mapa_textos_talento = {}
    mapa_textos_desarrollo = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=len(trabajos)) as fondo:
        futuros = {
            fondo.submit(redactor.redactar, **kwargs): (destino, pregunta)
            for destino, pregunta, kwargs in trabajos
        }
        for futuro in concurrent.futures.as_completed(futuros):
            destino, pregunta = futuros[futuro]
            texto = futuro.result()
            if destino == "familia":
                mapa_textos[pregunta] = texto
            elif destino == "talento":
                mapa_textos_talento[pregunta] = texto
            else:
                mapa_textos_desarrollo[pregunta] = texto

    # 3. Insertamos el texto en una copia del molde oficial
    molde_path = os.path.join(TEMPLATE_DIR, f"molde_{tipo}.pptx")
    if not os.path.exists(molde_path):
        return jsonify({"error": f"no se encontró el molde oficial: molde_{tipo}.pptx "
                                 f"(colóquelo en {TEMPLATE_DIR})"}), 500

    prs = Presentation(molde_path)
    slide = prs.slides[SLIDE_FAMILIA[tipo]]
    reporte = plantilla_pptx.llenar_respuestas(slide, mapa_textos)

    slide_talento = prs.slides[SLIDE_TALENTO[tipo]]
    # Letra más pequeña que en la slide de familia (size=12 por defecto):
    # este cuadro es más chico y trae 5 preguntas en vez de 4, así que con el
    # mismo tamaño el texto se desborda y se monta sobre las casillas vecinas.
    reporte_talento = plantilla_pptx.llenar_respuestas(slide_talento, mapa_textos_talento, size=9)
    reporte.update({f"[talento] {k}": v for k, v in reporte_talento.items()})

    # "Registro de observaciones al desarrollo infantil mensual" — solo en
    # llamada, y nunca cuando el participante es una mujer gestante (no se
    # hace registro de desarrollo infantil en ese caso).
    if tipo == "llamada" and tipo in SLIDE_DESARROLLO and not es_gestante:
        slide_desarrollo = prs.slides[SLIDE_DESARROLLO[tipo]]
        reporte_desarrollo = plantilla_pptx.llenar_respuestas(
            slide_desarrollo, mapa_textos_desarrollo, size=11)
        reporte.update({f"[desarrollo] {k}": v for k, v in reporte_desarrollo.items()})

    # 4. Entregamos el archivo
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    nombre_archivo = f"{datos['nombre'].strip().upper()} - voces - {uuid.uuid4().hex[:6]}.pptx"
    respuesta = send_file(
        buffer,
        as_attachment=True,
        download_name=nombre_archivo,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    respuesta.headers["X-Rayuela-Banda"] = banda["clave"]
    respuesta.headers["X-Rayuela-Reporte"] = str(reporte)
    return respuesta


# ─────────────────────────────────────────────────────────────────────────────
# MÓDULO DE PLANEACIÓN
# ─────────────────────────────────────────────────────────────────────────────

def _llenar_planeacion_en_pptx(prs, tipo, textos, objetos_paquete=""):
    """
    Escribe los textos de planeación en las diapositivas correspondientes
    del cuaderno (slides 0-1 para hogar, 0-3 para llamada).
    Usa el mismo motor espacial que voces (find_question / find_answer_table).
    """
    if tipo == "hogar":
        # Slide 0: INTENCIONALIDAD + Momento uno (Experiencias)
        slide0 = prs.slides[0]
        plantilla_pptx.llenar_respuestas(slide0, {
            "INTENCIONALIDAD": textos.get("intencionalidad", ""),
            "conectarnos": textos.get("experiencias_momento_uno", ""),
        }, size=11)

        # Slide 1: Momentos dos y tres + objetos paquete
        slide1 = prs.slides[1]
        mapa1 = {
            "construyendo juntos": textos.get("experiencias_momento_dos", ""),
            "comprometernos": textos.get("experiencias_momento_tres", ""),
        }
        if objetos_paquete:
            mapa1["paquete didáctico"] = objetos_paquete
        plantilla_pptx.llenar_respuestas(slide1, mapa1, size=11)

    else:  # llamada: DOS planeaciones independientes (slides 0-1 y 2-3)
        # Slide intro (0 y 2): INTENCIONALIDAD propia de cada planeación +
        # la "Experiencia pedagógica a promover con la familia" (texto fijo
        # de saludo + canción, igual para ambas salvo la canción).
        # Slide principal (1 y 3): la "Experiencia pedagógica principal",
        # completamente dinámica y propia de cada planeación.
        pares = [(0, 1, "1"), (2, 3, "2")]
        for slide_intro_idx, slide_principal_idx, n in pares:
            slide_intro = prs.slides[slide_intro_idx]
            plantilla_pptx.llenar_respuestas(slide_intro, {
                "INTENCIONALIDAD": textos.get(f"intencionalidad_{n}", ""),
                "Descripción de la experiencia": textos.get(f"intro_descripcion_{n}", ""),
                "Tiempo estimado y recursos": textos.get("intro_tiempo_recursos", ""),
            }, size=11)

            slide_principal = prs.slides[slide_principal_idx]
            mapa_principal = {
                "Descripción de la experiencia": textos.get(
                    f"experiencia_principal_descripcion_{n}", ""),
                "Tiempo estimado y recursos": textos.get(
                    f"experiencia_principal_tiempo_recursos_{n}", ""),
            }
            if objetos_paquete:
                mapa_principal["paquete didáctico"] = objetos_paquete
            plantilla_pptx.llenar_respuestas(slide_principal, mapa_principal, size=11)


@app.post("/api/planeacion")
def crear_planeacion():
    """
    Recibe los datos de planeación de Jimena, genera los textos con IA,
    guarda el registro en la BD y devuelve el PPTX de planeación para imprimir.

    Cuerpo JSON esperado para "hogar":
    {
      "nombre": "Hassan Melo Hueso",
      "fecha_nacimiento": "2024-09-09",
      "genero": "M",
      "tipo_cuaderno": "hogar",
      "actividad_principal": "exploración sensorial con texturas naturales",
      "nombre_ronda": "Los pollitos dicen",
      "link_ronda": "https://...",
      "objetos_paquete": "telas, semillas"   (opcional)
    }

    Cuerpo JSON esperado para "llamada" (dos planeaciones independientes):
    {
      "nombre": "Hassan Melo Hueso",
      "tipo_cuaderno": "llamada",
      "tipo_participante": "nino" | "nina" | "bebe" | "gestante",
      "fecha_nacimiento": "2024-09-09",   (no requerido si tipo_participante=="gestante")
      "genero": "M",                       (no requerido si tipo_participante=="gestante")
      "actividad_principal": "tema/experiencia de la planeación 1",
      "nombre_ronda": "canción 1", "link_ronda": "...",
      "actividad_principal_2": "tema/experiencia de la planeación 2",
      "nombre_ronda_2": "canción 2", "link_ronda_2": "...",
      "modalidad_acompanamiento": "Llamada telefónica"   (opcional),
      "objetos_paquete": "materiales disponibles en el hogar"   (opcional),
      "aspectos_fortalecer": "aspectos puntuales a fortalecer"   (opcional)
    }
    """
    datos = request.get_json(force=True)

    tipo = (datos.get("tipo_cuaderno") or "").strip().lower()
    if tipo not in ("hogar", "llamada"):
        return jsonify({"error": "tipo_cuaderno debe ser 'hogar' o 'llamada'"}), 400

    requeridos = ["nombre", "tipo_cuaderno", "actividad_principal"]
    if tipo == "llamada":
        requeridos += ["tipo_participante", "actividad_principal_2"]
    faltantes = [c for c in requeridos if not datos.get(c)]
    if faltantes:
        return jsonify({"error": f"faltan campos: {', '.join(faltantes)}"}), 400

    tipo_participante = (datos.get("tipo_participante") or "").strip().lower()
    if tipo == "llamada" and tipo_participante not in ("nino", "nina", "bebe", "gestante"):
        return jsonify({"error": "tipo_participante debe ser 'nino', 'nina', "
                                  "'bebe' o 'gestante'"}), 400

    es_gestante = tipo_participante == "gestante"
    if not es_gestante and (not datos.get("fecha_nacimiento") or not datos.get("genero")):
        return jsonify({"error": "faltan campos: fecha_nacimiento, genero "
                                  "(requeridos salvo mujer gestante)"}), 400

    genero = datos.get("genero", "M")

    if es_gestante:
        nacimiento = None
        banda_clave, banda_etiqueta = "gestante", "Mujer gestante"
        fecha_encuentro = datetime.now().date()
    else:
        try:
            nacimiento = datetime.strptime(datos["fecha_nacimiento"], "%Y-%m-%d").date()
        except ValueError:
            return jsonify({"error": "fecha_nacimiento inválida, use AAAA-MM-DD"}), 400
        meses = edades.edad_en_meses(nacimiento)
        banda = edades.banda_por_edad(meses)
        banda_clave, banda_etiqueta = banda["clave"], banda["etiqueta"]
        fecha_encuentro = nacimiento

    # Genera todos los textos de planeación con la API de Claude (en paralelo)
    if tipo == "hogar":
        textos = redactor.generar_textos_planeacion(
            actividad=datos["actividad_principal"],
            nombre_nino=datos["nombre"],
            banda_clave=banda_clave,
            banda_etiqueta=banda_etiqueta,
            nombre_ronda=datos.get("nombre_ronda", ""),
            link_ronda=datos.get("link_ronda", ""),
        )
    else:
        textos = redactor.generar_textos_planeacion_llamada(
            tema_1=datos["actividad_principal"],
            tema_2=datos["actividad_principal_2"],
            nombre_participante=datos["nombre"],
            tipo_participante=tipo_participante,
            genero=genero,
            banda_clave=banda_clave,
            banda_etiqueta=banda_etiqueta,
            nombre_cancion_1=datos.get("nombre_ronda", ""),
            link_cancion_1=datos.get("link_ronda", ""),
            nombre_cancion_2=datos.get("nombre_ronda_2", ""),
            link_cancion_2=datos.get("link_ronda_2", ""),
            modalidad_acompanamiento=datos.get("modalidad_acompanamiento", "Llamada telefónica"),
            materiales_disponibles=datos.get("objetos_paquete", ""),
            aspectos_fortalecer=datos.get("aspectos_fortalecer", ""),
        )

    # Abre el molde y llena las diapositivas de planeación
    molde_path = os.path.join(TEMPLATE_DIR, f"molde_{tipo}.pptx")
    if not os.path.exists(molde_path):
        return jsonify({"error": f"no se encontró el molde: molde_{tipo}.pptx"}), 500

    prs = Presentation(molde_path)
    _llenar_planeacion_en_pptx(prs, tipo, textos,
                                objetos_paquete=datos.get("objetos_paquete", ""))

    # Guarda el registro en la BD con estado 'pendiente_voces'
    planeacion_id = db.guardar_planeacion({
        "nombre_nino": datos["nombre"].strip(),
        "fecha_encuentro": fecha_encuentro.isoformat(),
        "genero": genero,
        "tipo_cuaderno": tipo,
        "tipo_participante": tipo_participante or None,
        "banda_clave": banda_clave,
        "banda_etiqueta": banda_etiqueta,
        "actividad_principal": datos["actividad_principal"],
        "actividad_principal_2": datos.get("actividad_principal_2", ""),
        "nombre_ronda": datos.get("nombre_ronda", ""),
        "link_ronda": datos.get("link_ronda", ""),
        "nombre_ronda_2": datos.get("nombre_ronda_2", ""),
        "link_ronda_2": datos.get("link_ronda_2", ""),
        "modalidad_acompanamiento": datos.get("modalidad_acompanamiento", ""),
        "objetos_paquete": datos.get("objetos_paquete", ""),
        "aspectos_fortalecer": datos.get("aspectos_fortalecer", ""),
        "textos_generados": textos,
    })

    # Devuelve el PPTX con las diapositivas de planeación listas para imprimir
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    nombre_archivo = (
        f"{datos['nombre'].strip().upper()} - planeacion - "
        f"{uuid.uuid4().hex[:6]}.pptx"
    )
    respuesta = send_file(
        buffer,
        as_attachment=True,
        download_name=nombre_archivo,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    respuesta.headers["X-Rayuela-Planeacion-Id"] = str(planeacion_id)
    return respuesta


@app.get("/api/planeaciones")
def listar_planeaciones():
    """Lista todas las planeaciones con estado 'pendiente_voces'."""
    try:
        pendientes = db.listar_pendientes()
        # Convierte fechas a string para que JSON las serialice correctamente
        for p in pendientes:
            for campo in ("fecha_encuentro", "fecha_creacion"):
                if p.get(campo) is not None:
                    p[campo] = str(p[campo])
        return jsonify(pendientes)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.get("/api/planeacion/<int:planeacion_id>")
def obtener_planeacion(planeacion_id):
    """Devuelve todos los datos de una planeación por su id."""
    plan = db.obtener_planeacion(planeacion_id)
    if plan is None:
        return jsonify({"error": "planeación no encontrada"}), 404
    for campo in ("fecha_encuentro", "fecha_creacion", "fecha_completado"):
        if plan.get(campo) is not None:
            plan[campo] = str(plan[campo])
    return jsonify(plan)


@app.post("/api/planeacion/<int:planeacion_id>/completar")
def completar_planeacion(planeacion_id):
    """
    Recibe las observaciones del encuentro real, genera el documento
    final con voces + planeación ya diligenciadas, y marca el registro
    como 'completado' en la BD.

    Cuerpo JSON esperado:
    {
      "observaciones": "le gustó mucho explorar las telas, se rio bastante..."
    }
    """
    datos = request.get_json(force=True)
    observaciones = (datos.get("observaciones") or "").strip()
    if not observaciones:
        return jsonify({"error": "el campo 'observaciones' es requerido"}), 400

    plan = db.obtener_planeacion(planeacion_id)
    if plan is None:
        return jsonify({"error": "planeación no encontrada"}), 404

    tipo = plan["tipo_cuaderno"]
    banda_clave = plan["banda_clave"]
    banda_etiqueta = plan["banda_etiqueta"]
    tipo_participante = (plan.get("tipo_participante") or "").strip().lower()
    es_gestante = tipo == "llamada" and tipo_participante == "gestante"
    textos_plan = plan.get("textos_generados") or {}

    # Reconstruye la materia prima para voces usando datos guardados + observaciones
    actividad_resumen = plan["actividad_principal"]
    if tipo == "llamada" and plan.get("actividad_principal_2"):
        actividad_resumen += f" y {plan['actividad_principal_2']}"
    materia_prima = (
        f"Actividad realizada: {actividad_resumen}\n"
        f"Lo que observé / lo que pasó: {observaciones}"
    )

    # Genera las voces (misma lógica que /api/generar-voces)
    preguntas = PREGUNTAS_HOGAR if tipo == "hogar" else PREGUNTAS_LLAMADA
    instrucciones = {
        preguntas[0]: "la respuesta de la familia en primera persona ('nosotros como "
                      "familia') sobre qué les gustó y qué no del encuentro/acompañamiento",
        preguntas[1]: "la respuesta de la familia en primera persona sobre sugerencias "
                      "para próximas experiencias",
        preguntas[2]: "la respuesta de la familia en primera persona sobre lo que "
                      "aprendieron de este encuentro/acompañamiento",
        preguntas[3]: "la respuesta de la familia en primera persona sobre los "
                      "compromisos para el próximo encuentro",
    }

    preguntas_talento = (PREGUNTAS_TALENTO_HOGAR if tipo == "hogar"
                         else PREGUNTAS_TALENTO_LLAMADA)
    instrucciones_talento = {
        preguntas_talento[0]: "mi reflexión como agente educativa, en primera persona "
                              "singular, sobre si se alcanzó la intencionalidad pedagógica",
        preguntas_talento[1]: "mi reflexión como agente educativa, en primera persona "
                              "singular, sobre cuál fue el mejor momento del encuentro",
    }
    if tipo == "hogar":
        instrucciones_talento.update({
            preguntas_talento[2]: "mi análisis sobre cómo participó la familia",
            preguntas_talento[3]: "mis recomendaciones para el próximo encuentro",
            preguntas_talento[4]: "mi valoración sobre cómo se aprovecharon los materiales",
        })
    else:
        instrucciones_talento[preguntas_talento[2]] = (
            "mi análisis sobre cómo se vinculó la familia en el acompañamiento a distancia"
        )
        instrucciones_talento[preguntas_talento[3]] = (
            "mis recomendaciones para los próximos acompañamientos a distancia"
        )

    preguntas_desarrollo = (PREGUNTAS_DESARROLLO_LLAMADA
                            if (tipo == "llamada" and not es_gestante) else [])
    instrucciones_desarrollo = {}
    if tipo == "llamada" and not es_gestante:
        primer_nombre = plan["nombre_nino"].strip().split()[0].title()
        instrucciones_desarrollo = {
            preguntas_desarrollo[0]: (
                f"una descripción narrativa, en tercera persona, nombrando a {primer_nombre} "
                f"por su nombre, sobre qué le gustó del encuentro y cómo participó"
            ),
            preguntas_desarrollo[1]: (
                f"una descripción narrativa, en tercera persona, sobre los intereses y "
                f"habilidades en desarrollo de {primer_nombre}"
            ),
            preguntas_desarrollo[2]: (
                f"mis recomendaciones profesionales, en primera persona singular, "
                f"para los próximos acompañamientos de {primer_nombre}"
            ),
        }

    # Lanza todas las llamadas a Claude en paralelo
    trabajos = []
    for pregunta in preguntas:
        trabajos.append(("familia", pregunta, dict(
            banda_clave=banda_clave, banda_etiqueta=banda_etiqueta,
            instruccion=instrucciones[pregunta], materia_prima=materia_prima,
            evitar_actividad=(tipo == "llamada"),
        )))
    for pregunta in preguntas_talento:
        trabajos.append(("talento", pregunta, dict(
            banda_clave=banda_clave, banda_etiqueta=banda_etiqueta,
            instruccion=instrucciones_talento[pregunta], materia_prima=materia_prima,
            perspectiva="talento_humano", max_palabras=70,
            evitar_actividad=(tipo == "llamada"),
        )))
    for pregunta in preguntas_desarrollo:
        trabajos.append(("desarrollo", pregunta, dict(
            banda_clave=banda_clave, banda_etiqueta=banda_etiqueta,
            instruccion=instrucciones_desarrollo[pregunta], materia_prima=materia_prima,
            perspectiva="desarrollo_infantil", max_palabras=90,
            evitar_actividad=True,
        )))

    mapa_textos, mapa_talento, mapa_desarrollo = {}, {}, {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=len(trabajos)) as fondo:
        futuros = {
            fondo.submit(redactor.redactar, **kwargs): (destino, pregunta)
            for destino, pregunta, kwargs in trabajos
        }
        for futuro in concurrent.futures.as_completed(futuros):
            destino, pregunta = futuros[futuro]
            if destino == "familia":
                mapa_textos[pregunta] = futuro.result()
            elif destino == "talento":
                mapa_talento[pregunta] = futuro.result()
            else:
                mapa_desarrollo[pregunta] = futuro.result()

    # Abre el molde y llena PRIMERO la planeación, LUEGO las voces
    molde_path = os.path.join(TEMPLATE_DIR, f"molde_{tipo}.pptx")
    if not os.path.exists(molde_path):
        return jsonify({"error": f"no se encontró el molde: molde_{tipo}.pptx"}), 500

    prs = Presentation(molde_path)

    # Planeación (reutiliza los textos ya guardados en la BD)
    _llenar_planeacion_en_pptx(prs, tipo, textos_plan,
                                objetos_paquete=plan.get("objetos_paquete", ""))

    # Voces
    slide_familia = prs.slides[SLIDE_FAMILIA[tipo]]
    plantilla_pptx.llenar_respuestas(slide_familia, mapa_textos)

    slide_talento = prs.slides[SLIDE_TALENTO[tipo]]
    plantilla_pptx.llenar_respuestas(slide_talento, mapa_talento, size=9)

    if tipo == "llamada" and tipo in SLIDE_DESARROLLO and not es_gestante:
        slide_desarrollo = prs.slides[SLIDE_DESARROLLO[tipo]]
        plantilla_pptx.llenar_respuestas(slide_desarrollo, mapa_desarrollo, size=11)

    # Marca como completado en la BD
    db.completar_planeacion(planeacion_id, observaciones)

    # Devuelve el cuaderno completo (planeación + voces)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    nombre_archivo = (
        f"{plan['nombre_nino'].upper()} - cuaderno completo - "
        f"{uuid.uuid4().hex[:6]}.pptx"
    )
    return send_file(
        buffer,
        as_attachment=True,
        download_name=nombre_archivo,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)), debug=True)
